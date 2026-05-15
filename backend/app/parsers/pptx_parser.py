from pptx import Presentation
from app.parsers.base import BaseParser, ParsedDocument, ParsedPage
from loguru import logger

class PPTXParser(BaseParser):
    def parse(self, file_path: str, file_id: str) -> ParsedDocument:
        logger.info(f"Parsing PPTX: {file_path}")
        prs = Presentation(file_path)
        pages = []
        
        for i, slide in enumerate(prs.slides):
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
            
            pages.append(ParsedPage(
                page_number=i + 1,
                text="\n".join(text_parts)
            ))
            
        return ParsedDocument(
            file_id=file_id,
            pages=pages,
            metadata={"total_slides": len(prs.slides)}
        )
